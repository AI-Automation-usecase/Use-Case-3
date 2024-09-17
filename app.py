import streamlit as st
import pandas as pd
import requests
import os
import zipfile
import fitz  # PyMuPDF
import re
from datetime import datetime
import numpy as np
import json
import streamlit.components.v1 as components
from io import BytesIO
from xlsxwriter import Workbook
from pptx import Presentation
from pptx.util import Inches, Pt  # Make sure Pt is imported
from pptx.enum.text import PP_ALIGN
import matplotlib.pyplot as plt
from collections.abc import Sequence
import time
from requests.exceptions import ChunkedEncodingError, ConnectionError, Timeout


# Set page configuration
st.set_page_config(layout="wide")

# CSS for header and footer
header_footer_css = """
    <style>
        .banner {
            position: relative;
            width: 100%;
            height: 150px;
            background-image: url('https://encrypted-tbn0.gstatic.com/images?q=tbn:ANd9GcQZJkPpcjFpNhfvdUcHdSytP1-tePz8v8X34Q&s');
            background-size: cover;
            background-position: center;
            color: transparent;
            text-align: center;
            line-height: 150px;
            transition: color 0.5s ease;
        }
        .banner:hover {
            color: white;
        }
        .logo {
            position: fixed;
            top: 0;
            right: 0;
            margin-right: 10px;
            margin-top: 10px;
            width: 100px;
            height: 150px;
        }
    </style>
    """
st.markdown(header_footer_css, unsafe_allow_html=True)
st.markdown("<div class='banner'></div>", unsafe_allow_html=True)

# Dashboard title and description
st.title("Financial Statements Dashboard")
st.image("https://img.freepik.com/premium-vector/business-statistics-financial-analytics-market-trend-analysis-vector-concept-illustration_92926-2486.jpg", caption="Trend Analysis and Data Insights", use_column_width=True)
st.sidebar.image("https://encrypted-tbn0.gstatic.com/images?q=tbn:ANd9GcRViRfutvNG9i9GtCPAC6qiwcK_uIOvKU0QP-zvFl3iMHKpUvAvStpetXH8o2AQ_fA4tBg&usqp=CAU", use_column_width=False)

# Radio buttons for navigation
page = st.sidebar.radio("Select Sheets", ["Balance Sheet", "Profit & Loss", "KPI"])

# Initialize session state variables
if 'uploaded_files' not in st.session_state:
    st.session_state['uploaded_files'] = []

# if 'balance_files_processed' not in st.session_state:
#     st.session_state['balance_files_processed'] = False
if 'balance_data_frames' not in st.session_state:
    st.session_state['balance_data_frames'] = []
if 'balance_results' not in st.session_state:
    st.session_state['balance_results'] = {}
if 'balance_fy_range' not in st.session_state:
    st.session_state['balance_fy_range'] = None
if 'balance_terms' not in st.session_state:
    st.session_state['balance_terms'] = []

# if 'balance_results_shown' not in st.session_state:
#     st.session_state['balance_results_shown'] = False
# if 'merged_df1' not in st.session_state:
#     st.session_state['merged_df1'] = {}

if 'include_forecast_bs' not in st.session_state:
    st.session_state['include_forecast_bs'] = False
if 'include_forecast_pl' not in st.session_state:
    st.session_state['include_forecast_pl'] = False

if 'pl_files_processed' not in st.session_state:
    st.session_state['pl_files_processed'] = False
if 'pl_data_frames' not in st.session_state:
    st.session_state['pl_data_frames'] = []
if 'pl_results_quarters' not in st.session_state:
    st.session_state['pl_results_quarters'] = {}
if 'pl_results_fy' not in st.session_state:
    st.session_state['pl_results_fy'] = {}
if 'pl_fy_range' not in st.session_state:
    st.session_state['pl_fy_range'] = None
if 'sorted_pl_quarters' not in st.session_state:
    st.session_state['sorted_pl_quarters'] = []
if 'sorted_pl_fys' not in st.session_state:
    st.session_state['sorted_pl_fys'] = []

if 'df_combined' not in st.session_state:
    st.session_state['df_combined'] = {}

if 'kpi_uploaded_files' not in st.session_state:
    st.session_state['kpi_uploaded_files'] = []
if 'kpi_file_processed' not in st.session_state:
    st.session_state['kpi_file_processed'] = False
if 'kpi_data' not in st.session_state:
    st.session_state['kpi_data'] = None
if 'kpi_results' not in st.session_state:
    st.session_state['kpi_results'] = None
if 'kpi_fy_columns' not in st.session_state:
    st.session_state['kpi_fy_columns'] = None

if 'downloaded_files' not in st.session_state:
    st.session_state['downloaded_files'] = {}



# Function to download files from the website only once
def download_files_once(company_symbol):
    if company_symbol in st.session_state['downloaded_files']:
        return st.session_state['downloaded_files'][company_symbol]
    else:
        company_name, all_files = download_files(company_symbol)
        st.session_state['downloaded_files'][company_symbol] = (company_name, all_files)
        return company_name, all_files


def extract_data_from_pdf_bs(pdf_file, heading_map, conversion_factor, pdf_filename):
    doc = fitz.open(stream=pdf_file, filetype="pdf")
    data = []
    date_columns = []
 
    total_pages = len(doc)
   
    consolidated_bs_pattern = r"Consolidated Balance Sheet.*"
 
    for page_num in range(total_pages):
        page_text = doc.load_page(page_num).get_text("text")
        #st.write(f"Debug: Reading page {page_num + 1} of {pdf_filename}")
        #st.write(f"Debug: Page content in {pdf_filename}: {page_text[:500]}...")  # Output the first 500 characters of the page
 
        # Look for the Consolidated Balance Sheet pattern
        match = re.search(consolidated_bs_pattern, page_text, re.IGNORECASE)
        if not match:
            #st.write(f"Debug: Consolidated Balance Sheet pattern not found on page {page_num + 1} of {pdf_filename}")
            continue
 
        # Check if this page has the required headings
        current_page_headings = {key for key, variations in heading_map.items() for heading in variations if re.search(rf"{re.escape(heading)}", page_text, re.IGNORECASE)}
        if current_page_headings == set(heading_map.keys()):
            date_columns = extract_date_columns(page_text)
            data = extract_values(page_text, heading_map, conversion_factor, pdf_filename)
            return data, date_columns
 
        # If not all headings are found, check the next page
        if page_num + 1 < total_pages:
            next_page_text = doc.load_page(page_num + 1).get_text("text")
            combined_text = page_text + "\n" + next_page_text
 
            current_combined_headings = {key for key, variations in heading_map.items() for heading in variations if re.search(rf"{re.escape(heading)}", combined_text, re.IGNORECASE)}
            if current_combined_headings == set(heading_map.keys()):
                date_columns = extract_date_columns(combined_text)
                data = extract_values(combined_text, heading_map, conversion_factor, pdf_filename)
                return data, date_columns
 
    #st.write(f"Debug: No relevant data found in file: {pdf_filename}")
    return None, None



def extract_date_columns(text):
    page_date_columns = re.findall(r"As at\s+(.*)\s+As at\s+(.*)", text, re.IGNORECASE)
    if not page_date_columns:
        page_date_columns = re.findall(r"(\b\w+ \d{1,2}, \d{4}\b)\s+(\b\w+ \d{1,2}, \d{4}\b)", text, re.IGNORECASE)
    if not page_date_columns:
        page_date_columns = re.findall(r"(\d{1,2} \b\w+ \d{4}\b)\s+(\d{1,2} \b\w+ \d{4}\b)", text, re.IGNORECASE)
    if not page_date_columns:
        page_date_columns = re.findall(r"As at March 31,\s+(\d{4})\s+(\d{4})", text, re.IGNORECASE)
        if page_date_columns:
            page_date_columns = [(f"March 31, {date.strip()}" for date in page_date_columns[0])]
    if not page_date_columns:
        page_date_columns = re.findall(r"As at March 31,\s+(\d{4})\s+(\d{4})", text, re.IGNORECASE)
        if page_date_columns:
            page_date_columns = [(f"March 31, {page_date_columns[0][0]}", f"March 31, {page_date_columns[0][1]}")]

    if not page_date_columns:
        return []

    date_list = []
    for date in page_date_columns[0]:
        for fmt in ("%B %d, %Y", "%d %B %Y"):
            try:
                date_obj = datetime.strptime(date.strip(), fmt)
                date_list.append(date_obj.strftime("%d-%m-%Y"))
                break
            except ValueError:
                continue
    return date_list

# Function to download files from the website with retry mechanism
def download_files(company_symbol, max_retries=3, retry_delay=5):
    url = f"https://www.nseindia.com/api/corp-info?symbol={company_symbol}&corpType=annualreport&market=cm"
    headers = {
        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/58.0.3029.110 Safari/537.3",
        "Accept": "application/json, text/plain, */*",
        "Referer": f"https://www.nseindia.com/get-quotes/equity?symbol={company_symbol}",
        "X-Requested-With": "XMLHttpRequest",
        "Connection": "keep-alive"
    }

    session = requests.Session()
    session.headers.update(headers)

    init_url = f"https://www.nseindia.com/get-quotes/equity?symbol={company_symbol}"
    session.get(init_url)

    response = session.get(url)

    if response.status_code == 200:
        data = response.json()
        if data:
            company_name_key = 'companyName' if 'companyName' in data[0] else 'company_name'
            company_name = data[0][company_name_key].replace(" ", "_")
            os.makedirs(company_name, exist_ok=True)

            file_names = [item['fileName'].split("/")[-1] for item in data[:4]]
            df_files = pd.DataFrame(file_names, columns=['FileName'])
            df_files['Exists'] = df_files['FileName'].apply(lambda x: os.path.isfile(os.path.join(company_name, x)))

            files_to_download = df_files[df_files['Exists'] == False]['FileName'].tolist()

            with st.spinner("Downloading files..."):
                for file_name in files_to_download:
                    file_url = next(item['fileName'] for item in data if item['fileName'].endswith(file_name))
                    file_extension = file_name.split('.')[-1].lower()

                    retries = 0
                    while retries < max_retries:
                        try:
                            file_response = session.get(file_url, stream=True)
                            if file_response.status_code == 200:
                                file_size = int(file_response.headers.get('Content-Length', 0))
                                downloaded_size = 0

                                file_path = os.path.join(company_name, file_name)
                                with open(file_path, 'wb') as file:
                                    for chunk in file_response.iter_content(1024 * 1024):  # 1MB chunks
                                        if chunk:
                                            file.write(chunk)
                                            downloaded_size += len(chunk)

                                # Check if the entire file has been downloaded
                                if downloaded_size < file_size:
                                    raise ChunkedEncodingError("Incomplete file download")

                                break  # If the download is successful, break the retry loop
                            else:
                                retries += 1
                                st.warning(f"Error {file_response.status_code}. Retrying {retries}/{max_retries}...")
                                time.sleep(retry_delay)
                        except (ChunkedEncodingError, ConnectionError, Timeout) as e:
                            retries += 1
                            st.warning(f"Error {e}. Retrying {retries}/{max_retries}...")
                            time.sleep(retry_delay)

                    if retries == max_retries:
                        st.error(f"Failed to download {file_name} after {max_retries} attempts.")
                        return None, []

                    # **Check if the file is a ZIP file before trying to open it**
                    if file_extension == 'zip':
                        try:
                            with zipfile.ZipFile(file_path, 'r') as zip_ref:
                                zip_ref.extractall(company_name)  # Extract ZIP contents
                        except zipfile.BadZipFile:
                            st.error(f"{file_name} is not a valid ZIP file or is corrupted.")
                            continue

            return company_name, df_files['FileName'].tolist()

    return None, []
# Function to load and process the data from the Excel file
def load_and_process_data(excel_file_path):
    excel_data = pd.ExcelFile(excel_file_path)
    processed_data = {}
    all_fy_columns = set()

    for sheet_name in excel_data.sheet_names:
        df = pd.read_excel(excel_file_path, sheet_name=sheet_name)
        df.set_index('Financial Year', inplace=True)

        # Collect all FY columns dynamically
        fy_columns = [col for col in df.columns if re.match(r'^FY \d{2}-\d{2}$', col)]
        all_fy_columns.update(fy_columns)
        
        processed_data[sheet_name] = df

    # Sort FY columns in ascending order based on the year part
    sorted_fy_columns = sorted(all_fy_columns, key=lambda x: int(x.split()[1].split('-')[0]))

    return processed_data, sorted_fy_columns


def extract_values(text, heading_map, conversion_factor, pdf_filename):
    with st.spinner(f"Fetching data from {pdf_filename}..."):
        lines = text.splitlines()
        data = []
        for key, variations in heading_map.items():
            for heading in variations:
                pattern = re.compile(rf"{re.escape(heading)}", re.IGNORECASE)
                for i, line in enumerate(lines):
                    if pattern.search(line.strip()):
                        try:
                            value_lines = lines[i+1:i+4]
                            numeric_values = []

                            for v in value_lines:
                                # Leave missing values as they are (None or empty)
                                v_clean = v.replace('\xa0', '').replace(',', '').strip()

                                # Check if the value is numeric
                                if re.match(r'^-?\d+(\.\d+)?(\s*\([ivx]+\))?$', v_clean, re.IGNORECASE) or re.match(r'^\d+(\.\d+)?$', v_clean):
                                    numeric_values.append(float(re.sub(r'\s*\([ivx]+\)', '', v_clean)))

                            if len(numeric_values) >= 2:
                                value1 = numeric_values[-2] * conversion_factor
                                value2 = numeric_values[-1] * conversion_factor
                            else:
                                value1, value2 = None, None  # Keep None for missing data

                            data.append([key, value1, value2])
                        except IndexError:
                            data.append([key, None, None])  # Keep None for missing data
                        except TypeError:
                            data.append([key, None, None])
                        break
        return data



def reformat_date_columns(date_columns):
    reformatted_columns = []
    for col in date_columns:
        try:
            date_obj = datetime.strptime(col, "%d-%m-%Y")
            reformatted_columns.append(date_obj.strftime("%d-%m-%Y"))
        except ValueError:
            reformatted_columns.append(col)
    return reformatted_columns

def extract_year(column_name):
    match = re.search(r'\d{4}', column_name)
    return int(match.group()) if match else 0

def calculate_financial_year(date_str):
    try:
        date_obj = datetime.strptime(date_str, "%d-%m-%Y")
        if date_obj.month > 3:
            fy_start = date_obj.year
            fy_end = date_obj.year + 1
        else:
            fy_start = date_obj.year - 1
            fy_end = date_obj.year
        return f"FY {str(fy_start)[-2:]}-{str(fy_end)[-2:]}"
    except ValueError:
        return None

# Function to display data for each heading across all companies based on FY values


def display_data_for_heading(processed_data, heading, fy_columns):
    all_data = []
    columns = ['Company'] + fy_columns

    for company, data in processed_data.items():
        if heading in data.index:
            row = [company] + [data.at[heading, fy] if fy in data.columns else None for fy in fy_columns]  # Keep None for missing values
        else:
            row = [company] + [None] * len(fy_columns)
        all_data.append(row)
    
    result_df = pd.DataFrame(all_data, columns=columns)
    
    # Sort the FY columns
    sorted_columns = ['Company'] + sort_fy_columns(result_df.columns[1:])
    result_df = result_df[sorted_columns]
    
    return result_df


def convert_lakhs_to_crores_value(value):
    if isinstance(value, (int, float)):
        return value / 100  # Convert lakhs to crores
    elif isinstance(value, str):
        cleaned_value = value.replace(',', '').strip()
        if cleaned_value.replace('.', '', 1).isdigit():
            return float(cleaned_value) / 100  # Convert lakhs to crores
    return value  # If the value is not a number or cannot be converted, return it as is

def convert_lakhs_to_crores(df):
    return df.applymap(convert_lakhs_to_crores_value)

def replace_missing_values(df):
    return df.replace({None: 0, 'None': 0, np.nan: 0})

# Function to read the file and extract the required sheet
def read_file(file, sheet_name):
    if file.type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
        df = pd.read_excel(file, sheet_name=None)
        if sheet_name in df.keys():
            df = pd.read_excel(file, sheet_name=sheet_name)
            # Convert values from Lakhs to Crores
            df = convert_lakhs_to_crores(df)  # Apply conversion to entire DataFrame
            return df, f"{file.name}_{sheet_name}"
        else:
            #st.error(f"The '{sheet_name}' sheet is not found in the uploaded file.")
            return None, file.name
    else:
        st.error("Unsupported file type. Please upload an Excel file.")
        return None, file.name

# Function to extract and sort financial years
def extract_and_sort_fy(df):
    fy_pattern = re.compile(r'FY \d{2}-\d{2}$')
    fys = [col for col in df.columns if fy_pattern.match(col)]
    fys.sort(key=lambda x: int(x.split()[1].split('-')[0]))
    return fys


def extract_and_sort_quarters(column_names):
    quarter_pattern = re.compile(r'Q[1-4] \d{2}-\d{2}$')
    quarters = [col for col in column_names if quarter_pattern.match(col)]
    
    # Sort quarters first by the financial year, then by the quarter within that year
    quarters.sort(key=lambda x: (int(x.split()[1].split('-')[0]), int(x.split()[1].split('-')[1]), x.split()[0]))
    
    return quarters

def process_extracted_data(final_data, extracted_data, extracted_date_columns, date_columns_set):
    temp_data = pd.DataFrame(extracted_data, columns=["Heading", "Value1", "Value2"])
    reformatted_date_columns = reformat_date_columns(extracted_date_columns)
 
    for col in reformatted_date_columns:
        if col not in final_data.columns:
            final_data[col] = None
            date_columns_set.add(col)
 
    for index, row in temp_data.iterrows():
        heading = row["Heading"]
        for i, col in enumerate(reformatted_date_columns):
            final_data.loc[final_data["Heading"] == heading, col] = row.iloc[i + 1]

# Function to process the dataframe and filter it based on search term
def process_dataframe(df, term, file_name):
    normalized_df = df.applymap(lambda x: x.strip().lower() if isinstance(x, str) else x)
    term_lower = term.lower()
    matching_rows = normalized_df[normalized_df.apply(lambda row: term_lower in row.values, axis=1)]
    if not matching_rows.empty:
        matching_rows.insert(0, 'File', file_name)
        return matching_rows  # Do not replace missing values with 0
    return pd.DataFrame()




# Function to generate the full range based on selected start and end
def generate_full_range(all_items, start, end):
    start_index = all_items.index(start)
    end_index = all_items.index(end) + 1  
    return all_items[start_index:end_index]



def forecast_and_plot(df, term, selected_items, col, add_legend=False, chart_type='line'):
    company_names = df['Company'].unique()
    all_series = []

    # Define constant colors for each company
    company_colors = {
        'sonata-software': '#1f77b4',   # Blue
        'TCS': '#ff7f0e',               # Orange
        'INFY': '#2ca02c',              # Green
        'HCLTECH': '#d62728',           # Red
        'LTIM': '#9467bd',              # Purple
        'WIPRO': '#8c564b',             # Brown
        'COFORGE': '#e377c2',           # Pink
        'PERSISTENT': '#7f7f7f',        # Gray
        'MPHASIS': '#bcbd22',           # Olive
        'ZENSARTECH': '#17becf'         # Cyan
    }

    # Terms for which the subtitle should not be displayed
    no_subtitle_terms = ["Basic (₹)", "EPS", "EBITDA Margin %", "Net Profit Margin %"]

    for company in company_names:
        company_data = df[df['Company'] == company]

        # Extracting and cleaning the data
        valid_values = company_data.iloc[0, 1:].values.tolist()
        filtered_values = [0 if pd.isna(value) else value for value in valid_values]
        filtered_values = np.array(filtered_values).astype(np.float64)
        filtered_items = list(selected_items[:len(filtered_values)])

        if len(filtered_values) > 0:
            series_data = {
                'name': company,
                'data': filtered_values.tolist(),
                'dashStyle': 'Solid',
                'dataLabels': {  # Enable data labels
                    'enabled': True,
                    'format': '{point.y:.2f}'  # Format to show two decimal places
                },
                'color': company_colors.get(company, '#000000')  # Assign color based on company, default to black
            }
            all_series.append(series_data)

    # Check if there's data to plot
    if not all_series:
        st.warning(f"No data available to plot for {term}.")
        return

    # Conditional subtitle based on the term
    if term in no_subtitle_terms:
        subtitle_text = ""
    else:
        subtitle_text = "Amount in ₹<br>₹ in Crores"

    highchart_html = f"""
    <script src="https://code.highcharts.com/highcharts.js"></script>
    <script src="https://code.highcharts.com/modules/annotations.js"></script>
    <div id="container_{term.replace(' ', '_')}" style="width:100%; height:500px;"></div>
    <script>
        Highcharts.chart('container_{term.replace(' ', '_')}', {{
            chart: {{
                zoomType: 'x',
                type: '{chart_type}'
            }},
            title: {{
                text: 'Visualization of {term}'
            }},
            subtitle: {{
                text: '{subtitle_text}',
                align: 'right',
                verticalAlign: 'top',
                style: {{
                    fontSize: '10px'
                }}
            }},
            xAxis: {{
                type: 'datetime',
                categories: {json.dumps(filtered_items)}
            }},
            yAxis: {{
                title: {{
                    text: 'Value'
                }}
            }},
            series: {json.dumps(all_series)},
            credits: {{
                enabled: false
            }},
            plotOptions: {{
                series: {{
                    dataLabels: {{
                        enabled: true
                    }}
                }}
            }}
        }});
    </script>
    """
    with col:
        components.html(highchart_html, height=500)




def display_consolidated_balance_results(results, fys, chart_type):
    valid_fys = sort_fy_columns(fys)  # Ensure FY columns are sorted
    terms = list(results.keys())
    
    include_forecast = st.session_state.get('include_forecast_bs', False)

    first_graph = True
    for term in terms:
        #st.subheader(f"Results for {term}")
        df = results[term].reindex(columns=['Company'] + valid_fys)
        st.dataframe(df)

        # Plotting
        if not df.empty:
            cols = st.columns(1)
            forecast_and_plot(df, term, valid_fys, cols[0], add_legend=first_graph, chart_type=chart_type)
            first_graph = False


def aggregate_data_single_row_merged(search_terms, data_frames_info, selected_fys):
    company_name = "sonata-software"
    aggregate_results = {}

    for term in search_terms:
        row_data = {'Company': company_name}
        for fy in selected_fys:
            values = []
            for df_info in data_frames_info:
                df, file_name = df_info
                if df is not None:
                    processed_df = process_dataframe(df, term, file_name)
                    if fy in processed_df.columns and not processed_df.empty:
                        value = processed_df[fy].iloc[0] if pd.notna(processed_df[fy].iloc[0]) else 0
                        if term == 'Basic (₹)':  # Apply multiplication for Sonata Software data
                            value *= 100  # Multiply by 100
                        values.append(value)
            row_data[fy] = values[0] if values else 0
        aggregate_results[term] = pd.DataFrame([row_data])

    return aggregate_results


# Ensure you only download and process files once
def initialize_download_and_process():
    # Path for the output Excel file
    output_excel_file_path = "peers_data_bs_demo.xlsx"

    # Check if the data has already been processed
    if 'processed_data' not in st.session_state or 'fy_columns' not in st.session_state:
        with st.spinner("Downloading and processing files..."):
            # Define the companies and their respective units
            companies = {
                "TCS": "crores",
                "WIPRO": "million",
                "MPHASIS": "million",
                "INFY": "crores",
                "LTIM": "million",
                "ZENSARTECH": "million",
                "HCLTECH": "crores",
                "COFORGE": "million"
            }

            # Define conversion factors for different units
            conversion_factors = {
                "crores": 1,          # Already in crores
                "million": 0.1        # Convert from million to crores
            }

            # Create an Excel writer to save processed data
            with pd.ExcelWriter(output_excel_file_path, engine='xlsxwriter') as writer:
                for company, unit in companies.items():
                    conversion_factor = conversion_factors[unit]
                    company_name, all_files = download_files_once(company)

                    if company_name:
                        # Debug output for the number of files and their names
                        #st.write(f"Debug: {len(all_files)} files to be processed for company: {company}")
                        #st.write(f"Debug: Files for company {company}: {all_files}")

                        # Map of required headings to their possible variations
                        heading_map = {
                            "Total non-current assets": ["Total non-current assets", "Total Non-current assets", "Total Non-Current Assets","Total non-current assets","Total non-current assets","Total Non-current assets"],
                            "Total current assets": ["Total current assets", "Total Current Assets","otal current assets"],
                            "Total assets": ["TOTAL ASSETS", "Total assets", "Total Assets","TOTAL ASSETS","TOTAL ASSETS","Total assets"],
                            "Total Equity": ["Total equity", "Total Equity", "TOTAL EQUITY","TOTAL EQUITY","Total equity","Total Equity"],
                            "Total non-current liabilities": ["Total non-current liabilities", "Total non- current liabilities", "Total Non-Current Liabilities","Total non-current liabilities","Total non-current liabilities", "Total non- current liabilities"],
                            "Total current liabilities": ["Total current liabilities", "Total Current Liabilities","Total current liabilities","Total current liabilities"],
                            "Total equity and liabilities": ["Total equity and liabilities", "Total Equity and Liabilities", "TOTAL EQUITY AND LIABILITIES","TOTAL EQUITY AND LIABILITIES","Total equity and liabilities"]
                                
                        }  

                        # Initialize the final DataFrame with the required headings
                        final_data = pd.DataFrame(columns=["Heading"] + list(heading_map.keys()))
                        final_data["Heading"] = list(heading_map.keys())

                        processed_files = set()
                        date_columns_set = set()

                        for file_name in all_files:
                            file_path = os.path.join(company_name, file_name)

                            if file_name.endswith(".zip"):
                                # Handle ZIP file
                                with zipfile.ZipFile(file_path, 'r') as zip_ref:
                                    pdf_file_names = [file for file in zip_ref.namelist() if file.endswith('.pdf')]
                                    for pdf_file_name in pdf_file_names:
                                        if pdf_file_name not in processed_files:
                                            with zip_ref.open(pdf_file_name) as pdf_file:
                                                extracted_data, extracted_date_columns = extract_data_from_pdf_bs(pdf_file.read(), heading_map, conversion_factor, pdf_file_name)
                                                if extracted_data and extracted_date_columns:
                                                    process_extracted_data(final_data, extracted_data, extracted_date_columns, date_columns_set)
                                                    processed_files.add(pdf_file_name)
                            elif file_name.endswith(".pdf"):
                                # Handle PDF file
                                if file_name not in processed_files:
                                    #st.write(f"Debug: Attempting to process file: {file_name}")
                                    with open(file_path, 'rb') as pdf_file:
                                        extracted_data, extracted_date_columns = extract_data_from_pdf_bs(pdf_file.read(), heading_map, conversion_factor, file_name)
                                        if extracted_data and extracted_date_columns:
                                            process_extracted_data(final_data, extracted_data, extracted_date_columns, date_columns_set)
                                            processed_files.add(file_name)
                                        else:
                                            st.write(f"Debug: No data extracted from file: {file_name}")

                        # Remove the initialized columns used for structure
                        final_data = final_data.drop(columns=list(heading_map.keys()), axis=1, errors='ignore')

                        # Sort the columns based on the year extracted from the column names
                        columns_sorted = ["Heading"] + sorted(list(date_columns_set), key=extract_year)
                        final_data = final_data[columns_sorted]

                        # Add a new row at the top for financial year values
                        financial_year_row = ["Financial Year"] + [calculate_financial_year(col) for col in columns_sorted[1:]]
                        final_data.loc[-1] = financial_year_row  # Add the row at the start
                        final_data.index = final_data.index + 1  # Shift index by 1
                        final_data = final_data.sort_index()  # Sort the index

                        # Write the DataFrame to an Excel sheet
                        final_data.to_excel(writer, sheet_name=company, index=False, header=False)

            # Load the processed Excel file into session state
            processed_data, fy_columns = load_and_process_data(output_excel_file_path)
            st.session_state['processed_data'] = processed_data
            st.session_state['fy_columns'] = fy_columns

# Ensure the initialization is done only once
initialize_download_and_process()


# Function for Balance Sheet Page
def filter_none_rows(df):
    fy_columns = [col for col in df.columns if col.startswith('FY')]
    df_filtered = df.dropna(subset=fy_columns, how='all')  # Remove rows where all FY columns are None
    return df_filtered


def sort_fy_columns(fy_columns):
    # Filter out any columns that do not match the FY format
    fy_columns_filtered = [fy for fy in fy_columns if re.match(r'^FY \d{2}-\d{2}$', fy)]
    
    # Add a debug statement to see what columns are being processed
    #st.write("Debug: Columns being processed for sorting:", fy_columns_filtered)
    
    # Proceed with sorting only the filtered columns
    fy_years = [int(fy.split()[1].split('-')[0]) for fy in fy_columns_filtered]
    sorted_indices = sorted(range(len(fy_years)), key=lambda k: fy_years[k])
    
    return [fy_columns_filtered[i] for i in sorted_indices]


# Function to generate a graph using matplotlib (only for PowerPoint)
def save_matplotlib_graph(df, term):
    # Define constant colors for each company (same as used in Highcharts)
    company_colors = {
        'sonata-software': '#1f77b4',   # Blue
        'TCS': '#ff7f0e',               # Orange
        'INFY': '#2ca02c',              # Green
        'HCLTECH': '#d62728',           # Red
        'LTIM': '#9467bd',              # Purple
        'WIPRO': '#8c564b',             # Brown
        'COFORGE': '#e377c2',           # Pink
        'PERSISTENT': '#7f7f7f',        # Gray
        'MPHASIS': '#bcbd22',           # Olive
        'ZENSARTECH': '#17becf'         # Cyan
    }

    # Terms that should not have the subtitle
    exclude_terms = ["Basic (₹)", "EPS", "EBITDA Margin %", "Net Profit Margin %"]

    fig, ax = plt.subplots(figsize=(12, 6))  # Increase the figure size to give more room

    # Plot data with constant colors for each company
    for company in df['Company'].unique():
        company_data = df[df['Company'] == company].set_index('Company').T
        ax.plot(company_data.index, company_data.values, label=company, 
                marker='o', color=company_colors.get(company, '#000000'))  # Default to black if company is not found

    # Set the title, x-axis, and y-axis labels
    ax.set_title(term, fontsize=16)
    ax.set_ylabel('Value', fontsize=12)
    ax.set_xlabel('Fiscal Year / Quarter', fontsize=12)

    # Add subtitle in the top-right corner, only for specific terms
    if term not in exclude_terms:
        subtitle_text = "Amount in ₹\n₹ in Crores"
        # Position subtitle on the top-right corner of the plot
        ax.text(1, 1.02, subtitle_text, transform=ax.transAxes, ha='right', fontsize=10, color='gray')

    # Rotate the x-axis labels for better readability
    plt.xticks(rotation=45, fontsize=10)

    # Move the legend to the bottom, increase number of columns, and reduce font size
    ax.legend(loc='upper center', bbox_to_anchor=(0.5, -0.2), ncol=6, fontsize=8, frameon=False)

    # Adjust layout to fit the legend at the bottom
    plt.tight_layout(rect=[0, 0, 1, 1])

    # Save the graph as an image to use in PowerPoint
    image_stream = BytesIO()
    fig.savefig(image_stream, format='png')
    image_stream.seek(0)

    return image_stream


# Function to generate and return PowerPoint presentation
def generate_ppt(term_data_dict, title):
    prs = Presentation()

    for term, df in term_data_dict.items():
        # Remove duplicates from the dataframe
        df = df.drop_duplicates(subset='Company', keep='first')

        # Filter out rows that are entirely NaN or empty
        df = df.dropna(how='all')

        # Recalculate rows and columns after removing duplicates and empty rows
        rows, cols = df.shape

        # Add slide for table
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        # Add title to slide
        title_shape = slide.shapes.title
        title_shape.text = f"{term}"

        # Add table to slide (using dynamic row count based on filtered DataFrame)
        table = slide.shapes.add_table(rows + 1, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(4)).table

        # Set column names in the first row
        for i, col_name in enumerate(df.columns):
            if i < cols:  # Ensure we're within column bounds
                table.cell(0, i).text = col_name
                # Decrease the font size of the header
                for paragraph in table.cell(0, i).text_frame.paragraphs:
                    paragraph.font.size = Pt(10)  # Set font size for column headers

        # Set data in the table (formatted to 3 decimal places)
        current_row = 1  # Start at the first data row
        for i, row in df.iterrows():
            # Check if the entire row is empty (all values are NaN)
            if row.isnull().all():
                continue  # Skip this row if it's entirely empty

            # Add only non-empty rows to the table
            for j, value in enumerate(row):
                if j < cols:  # Ensure we're within column bounds
                    if isinstance(value, (int, float)):
                        table.cell(current_row, j).text = f"{value:.3f}"  # Limit to 3 decimal places
                    else:
                        table.cell(current_row, j).text = str(value)

                    # Decrease the font size of the data
                    for paragraph in table.cell(current_row, j).text_frame.paragraphs:
                        paragraph.font.size = Pt(9)  # Set font size for data

            # Move to the next row in the table for non-empty data
            current_row += 1

        # Add slide for graph
        slide = prs.slides.add_slide(prs.slide_layouts[5])

        # Generate and save graph image
        image_stream = save_matplotlib_graph(df, term)

        # Insert the saved graph image into the slide
        slide.shapes.add_picture(image_stream, Inches(1), Inches(2), Inches(8), Inches(5.5))

    # Save presentation to a BytesIO object
    ppt_io = BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)

    return ppt_io


# Updated balance_sheet_page function



def balance_sheet_page():
    term_data_dict = {}

    st.subheader("Balance Sheet")

    # Track peer comparison state globally if not already set
    if 'peer_comparison_enabled' not in st.session_state:
        st.session_state['peer_comparison_enabled'] = False

    peer_comparison_enabled = st.checkbox("Peer Comparison", value=st.session_state['peer_comparison_enabled'])
    st.session_state['peer_comparison_enabled'] = peer_comparison_enabled

    # Handle file uploads
    uploaded_files = st.sidebar.file_uploader("Choose Excel files for BS & PL", type='xlsx', accept_multiple_files=True, key='balance_files')

    if uploaded_files:
        if 'uploaded_files' not in st.session_state:
            st.session_state['uploaded_files'] = []
        # Add newly uploaded files to session state
        for file in uploaded_files:
            if file not in st.session_state['uploaded_files']:
                st.session_state['uploaded_files'].append(file)
                st.session_state['balance_files_processed'] = False  # Mark as not processed

    if 'merged_df1' not in st.session_state:
        st.session_state['merged_df1'] = {}

    # Display the uploaded files and allow removal
    if st.session_state['uploaded_files']:
        with st.sidebar.expander("View Uploaded Files"):
            for i, file in enumerate(st.session_state['uploaded_files']):
                col1, col2 = st.columns([8, 2])
                col1.write(file.name)
                if col2.button("X", key=f"remove_balance_{i}"):
                    st.session_state['uploaded_files'].pop(i)
                    st.session_state['balance_files_processed'] = False  # Mark as not processed
                    st.rerun()  # Refresh the page automatically

    # Automatically process files when they are uploaded or deleted
    if not st.session_state.get('balance_files_processed') and st.session_state['uploaded_files']:
        with st.spinner("Processing files..."):
            data_frames_info = []
            for file in st.session_state['uploaded_files']:
                try:
                    df_info = read_file(file, 'BS_Y')  # Assuming read_file is a custom function you defined
                    if df_info[0] is None:
                        st.error(f"The file {file.name} does not contain the required sheet 'BS_Y'. Please upload a file with the correct sheet name.")
                        return
                    data_frames_info.append(df_info)
                except Exception as e:
                    st.error(f"Error processing {file.name}: {e}")
                    return

            st.session_state['balance_data_frames'] = data_frames_info
            st.session_state['balance_files_processed'] = True
            st.rerun()  # Refresh the page automatically

    # Check if data is processed correctly
    if st.session_state.get('balance_files_processed'):
        fys = set()
        for df_info in st.session_state['balance_data_frames']:
            df, _ = df_info
            fys.update(extract_and_sort_fy(df))
        fys = sorted(fys, key=lambda x: int(x.split()[1].split('-')[0]))

        if fys and st.session_state['balance_fy_range'] is None:
            st.session_state['balance_fy_range'] = (fys[0], fys[-1])

        selected_fys = []
        if peer_comparison_enabled:
            selected_fys = fys
        else:
            selected_fy_range = st.select_slider("Select Financial Year Range", options=fys, value=st.session_state['balance_fy_range'], key='balance_fy_range')
            if selected_fy_range:
                selected_fys = generate_full_range(fys, selected_fy_range[0], selected_fy_range[1])

        search_terms = [
            "Total non-current assets", "Total current assets", "Total assets", 
            "Total Equity", "Total non-current liabilities", "Total current liabilities", 
            "Total equity and liabilities"
        ]
        selected_search_terms = st.multiselect("Select Search Terms", options=search_terms, default=search_terms, key='balance_terms')

        # Only show the company selection when peer comparison is enabled
        if peer_comparison_enabled:
            companies = ['sonata-software', 'TCS', 'INFY', 'HCLTECH', 'LTIM', 'WIPRO', 'COFORGE', 'MPHASIS', 'ZENSARTECH']
            selected_companies = st.multiselect("Select Companies for Comparison", options=companies, default=companies, key='selected_companies')

        # Track if "Show BS Results" button is clicked or if results have already been generated
        if 'balance_show_results_clicked' not in st.session_state:
            st.session_state['balance_show_results_clicked'] = False

        # If the button is clicked, update the flag
        if st.button('Show BS Results', key='balance_show_results') or st.session_state['balance_show_results_clicked']:
            st.session_state['balance_show_results_clicked'] = True  # Store the result in session state

            # Process and display results
            results = aggregate_data_single_row_merged(selected_search_terms, st.session_state['balance_data_frames'], selected_fys)
            st.session_state['balance_results'] = results
            st.session_state['sorted_balance_fys'] = selected_fys

        display_results = st.session_state['balance_show_results_clicked']

        # Show the results if the button has been clicked or state indicates to show results
        if display_results:
            if peer_comparison_enabled:
                #st.subheader("Peer Comparison")
                for term in selected_search_terms:
                    # Extract peer data and Sonata data for the term
                    peer_data_df = display_data_for_heading(st.session_state['processed_data'], term, st.session_state['fy_columns'])
                    sonata_data_df = st.session_state['balance_results'][term]

                    # Filter the peer data to only show selected companies
                    peer_data_df = peer_data_df[peer_data_df['Company'].isin(selected_companies)]

                    # Extract ZENSARTECH data from peer_data_df
                    zensartech_data = peer_data_df[peer_data_df['Company'] == 'ZENSARTECH']

                    # Check if ZENSARTECH is already in sonata_data_df to prevent duplication
                    if not zensartech_data.empty and 'ZENSARTECH' not in sonata_data_df['Company'].values:
                        # Add ZENSARTECH data to sonata_data_df
                        sonata_data_df = pd.concat([sonata_data_df, zensartech_data])

                    # Extract the columns from sonata_data_df that match the columns in peer_data_df
                    sonata_columns = [col for col in peer_data_df.columns if col in sonata_data_df.columns and col != 'Company']

                    # Initialize missing columns in sonata_data_df as zeros if they are in peer_data_df but not in sonata_data_df
                    missing_columns = [col for col in peer_data_df.columns if col not in sonata_data_df.columns and col != 'Company']
                    for col in missing_columns:
                        sonata_data_df[col] = 0

                    # Ensure 'Company' is added only once at the start
                    sonata_data_df = sonata_data_df[['Company'] + sonata_columns + missing_columns]

                    # Merge the data
                    merged_df1 = pd.concat([sonata_data_df, peer_data_df], ignore_index=True)

                    # Drop duplicates to ensure ZENSARTECH or any other company doesn't appear twice
                    merged_df1 = merged_df1.drop_duplicates(subset='Company', keep='first')

                    if "ZENSARTECH" in merged_df1['Company'].values:
                        zensartech_data = merged_df1[merged_df1['Company'] == "ZENSARTECH"]
                        for fy in ['FY 19-20', 'FY 20-21']:
                            if fy in zensartech_data.columns:
                                original_value = zensartech_data[fy].values
                                adjusted_value = original_value * 0.1
                                zensartech_data[fy] = adjusted_value
                        merged_df1.update(zensartech_data)

                    # Filter out rows with only None values in FY columns
                    merged_df1 = filter_none_rows(merged_df1)

                    st.session_state['merged_df1'][term] = merged_df1
                    term_data_dict[term] = merged_df1
                    st.dataframe(merged_df1)

                    # Plot the merged data
                    cols = st.columns(1)
                    forecast_and_plot(merged_df1, term, merged_df1.columns[1:], cols[0], add_legend=False, chart_type='line')
                    st.markdown("**Note:** Data indicating zero may imply that the data is not present or not extracted due to improper headings.")

            else:
                display_consolidated_balance_results(st.session_state['balance_results'], st.session_state['sorted_balance_fys'], chart_type='area')

            if peer_comparison_enabled and term_data_dict:
                ppt_io = generate_ppt(term_data_dict, "BL Presentation")
                st.download_button(
                    label="Download BL PowerPoint",
                    data=ppt_io,
                    file_name="balance_sheet_visualization.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )




def calculate_quarter(date):
    month = date.month
    year = date.year
    if month in [1, 2, 3]:
        quarter = 'Q4'
        financial_year = f'{str(year-1)[-2:]}-{str(year)[-2:]}'
    elif month in [4, 5, 6]:
        quarter = 'Q1'
        financial_year = f'{str(year)[-2:]}-{str(year+1)[-2:]}'
    elif month in [7, 8, 9]:
        quarter = 'Q2'
        financial_year = f'{str(year)[-2:]}-{str(year+1)[-2:]}'
    else:
        quarter = 'Q3'
        financial_year = f'{str(year)[-2:]}-{str(year+1)[-2:]}'
    return f'{quarter} {financial_year}'

def initialize_session(symbol):
    headers = {
        'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.3',
        'Accept-Language': 'en-US,en;q=0.9',
        'Accept-Encoding': 'gzip, deflate, br',
        'Connection': 'keep-alive',
        'Referer': f'https://www.nseindia.com/get-quotes/equity?symbol={symbol}',
        'DNT': '1',
        'sec-fetch-dest': 'empty',
        'sec-fetch-mode': 'cors',
        'sec-fetch-site': 'same-origin'
    }

    session = requests.Session()
    session.headers.update(headers)

    init_url = f"https://www.nseindia.com/get-quotes/equity?symbol={symbol}"
    response = session.get(init_url)

    if response.status_code == 200:
        return session
    else:
        st.error("Failed to initialize session.")
        return None

def fetch_nse_financial_results(symbol, session):
    url = f"https://www.nseindia.com/api/corporates-financial-results?index=equities&symbol={symbol}&period=Quarterly"
    response = session.get(url)
    
    if response.status_code == 200:
        try:
            return response.json()
        except ValueError:
            st.error("Response content is not valid JSON.")
            st.write(response.text)
            return None
    else:
        st.error(f"Failed to fetch data for {symbol}. HTTP Status Code: {response.status_code}")
        st.write(response.text)
        return None

def create_financial_result_links(data):
    links = []
    
    for item in data:
        symbol = item.get("symbol", "")
        from_date = item.get("fromDate", "")
        to_date = item.get("toDate", "")
        seq_number = item.get("seqNumber", "")
        quarter = get_quarter_from_date(to_date)
        
        if symbol and from_date and to_date and quarter and seq_number:
            url = (
                f"https://www.nseindia.com/api/corporates-financial-results-data?"
                f"index=equities&params={from_date}{to_date}{quarter}ANNCNE{symbol}&seq_id={seq_number}"
                "&industry=-&frOldNewFlag=N&ind=N&format=New"
            )
            links.append(url)
    
    return links

def get_quarter_from_date(to_date):
    try:
        date_obj = datetime.strptime(to_date, '%d-%b-%Y')
        return calculate_quarter(date_obj)
    except ValueError:
        return "N/A"

def filter_data_by_year_and_consolidated(data, max_results=8):
    filtered_data = []
    seen_xbrl = set()
    
    for item in data:
        if len(filtered_data) >= max_results:
            break
        try:
            consolidated_status = item.get("consolidated", "")
            xbrl_url = item.get("xbrl", "")

            if consolidated_status == "Consolidated" and xbrl_url not in seen_xbrl:
                filtered_data.append(item)
                seen_xbrl.add(xbrl_url)
        except ValueError:
            continue
    
    return filtered_data

def fetch_financial_result_details(url, session):
    response = session.get(url)
    
    if response.status_code == 200:
        try:
            return response.json()
        except ValueError:
            st.error("Response content is not valid JSON.")
            st.write(response.text)
            return None
    elif response.status_code == 401:
        st.error("Unauthorized access. Ensure that session is properly initialized and headers are correct.")
    else:
        st.error(f"Failed to fetch data. HTTP Status Code: {response.status_code}")
        st.write(response.text)
        return None

def extract_required_fields(details, quarter, symbol):
    extracted_data = {
        "Company": symbol,
        "Quarter": quarter,
        "Total Income": details.get("re_total_inc"),
        "Total expenses": details.get("re_oth_tot_exp"),
        "Net tax expense": details.get("re_tax"),
        "Profit for the year": details.get("re_con_pro_loss"),
        "Revenue from operations": details.get("re_net_sale"),
        "Finance costs": details.get("re_int_new"),
        "Depreciation and amortization expense": details.get("re_depr_und_exp"),
        "Profit before exceptional item and tax": details.get("re_pro_bef_int_n_excep"),
        "Basic (₹)": details.get("re_basic_eps_for_cont_dic_opr"),
    }
    return extracted_data



def pivot_data_to_format(results_quarters, selected_search_terms_quarters):
    formatted_data = {}
    
    for term in selected_search_terms_quarters:
        if term in results_quarters:
            df = results_quarters[term]
            
            # Convert the 'Value' column to numeric, forcing non-numeric values to NaN
            df['Value'] = pd.to_numeric(df['Value'], errors='coerce')
            
            # Pivot the DataFrame to have Companies as rows and Quarters as columns
            df_pivot = df.pivot_table(index='Company', columns='Quarter', values='Value', aggfunc='sum')
            
            # Ensure the columns (quarters) are sorted correctly
            sorted_quarters = extract_and_sort_quarters(df_pivot.columns)
            df_pivot = df_pivot[sorted_quarters]
            
            formatted_data[term] = df_pivot
    
    return formatted_data



def profit_loss_page():
    # Centralize file upload handling and data processing
    term_data_dict = {}  # Dictionary to hold data for each term
    st.subheader("Profit & Loss")

    uploaded_files = st.sidebar.file_uploader("Choose Excel files for BS & PL", type='xlsx', accept_multiple_files=True, key='pl_files')

    if uploaded_files:
        if 'uploaded_files' not in st.session_state:
            st.session_state['uploaded_files'] = []
        # Add newly uploaded files to session state
        for file in uploaded_files:
            if file not in st.session_state['uploaded_files']:
                st.session_state['uploaded_files'].append(file)
                st.session_state['pl_files_processed'] = False  # Mark as not processed

    # Display the uploaded files and allow removal
    if st.session_state['uploaded_files']:
        with st.sidebar.expander("View Uploaded Files"):
            for i, file in enumerate(st.session_state['uploaded_files']):
                col1, col2 = st.columns([8, 2])
                col1.write(file.name)
                if col2.button("X", key=f"remove_pl_{i}"):
                    st.session_state['uploaded_files'].pop(i)
                    st.session_state['pl_files_processed'] = False  # Mark as not processed
                    st.rerun()

    # Automatically process files when they are uploaded, removed, or when the page loads
    if not st.session_state.get('pl_files_processed') and st.session_state['uploaded_files']:
        with st.spinner("Processing files..."):
            data_frames_info = []
            sheet_present = True
            for file in st.session_state['uploaded_files']:
                try:
                    df_info = read_file(file, 'PL_Q')
                    if df_info[0] is None:
                        st.error(f"The file {file.name} does not contain the required sheet 'PL_Q'. Please upload a file with the correct sheet name.")
                        sheet_present = False
                        break
                    data_frames_info.append(df_info)
                except Exception as e:
                    st.error(f"Error processing {file.name}: {e}")
                    return

            if sheet_present:
                st.session_state['pl_data_frames'] = data_frames_info
                st.session_state['pl_files_processed'] = True
                st.rerun()

    # Check if data is processed correctly
    if st.session_state.get('pl_files_processed'):
        # Debugging: Print the processed data to check if it's correct
        #st.write("Processed PL Data Frames:", st.session_state['pl_data_frames'])

        # Fetch and store NSE data if not already done
        if 'nse_data_processed' not in st.session_state:
            with st.spinner("Fetching NSE data..."):
                fetch_and_store_nse_data()

        quarters = set()
        for df_info in st.session_state['pl_data_frames']:
            df, _ = df_info
            quarters.update(extract_and_sort_quarters(df.columns))

        quarters = sorted(quarters, key=lambda x: (int(x.split()[1].split('-')[0]), x.split()[0][1]))

        if quarters and st.session_state['pl_fy_range'] is None:
            st.session_state['pl_fy_range'] = (quarters[0], quarters[-1])

        peer_comparison_enabled = st.checkbox("Peer Comparison")

        # Only show the company selection when peer comparison is enabled
        if peer_comparison_enabled:
            companies = ['TCS', 'INFY', 'HCLTECH', 'LTIM', 'WIPRO', 'COFORGE', 'PERSISTENT', 'MPHASIS', 'ZENSARTECH']
            selected_companies = st.multiselect("Select Companies for Comparison", options=companies, default=companies, key='selected_companies_pl')

        # Only display the quarter range slider when peer comparison is disabled
        selected_quarters = []
        if peer_comparison_enabled:
            selected_quarters = quarters  # Use all available quarters when Peer Comparison is enabled
        else:
            selected_quarter_range = st.select_slider(
                "Select Quarter Range",
                options=quarters,
                value=st.session_state['pl_fy_range'],
                key='pl_fy_range',
                disabled=peer_comparison_enabled  # Disable when peer comparison is enabled
            )
            if selected_quarter_range:
                selected_quarters = generate_full_range(quarters, selected_quarter_range[0], selected_quarter_range[1])

        all_search_terms_quarters = [
            "Total Income", "Total expenses", "Net tax expense", "Profit for the year",
            "Revenue from operations", "Finance costs", "Depreciation and amortization expense",
            "Profit before exceptional item and tax", "Basic (₹)"
        ]

        selected_search_terms_quarters = st.multiselect(
            "Select Search Terms for Quarters",
            options=all_search_terms_quarters,
            default=all_search_terms_quarters
        )

        # Initialize df_combined in session state if not already present
        if 'df_combined' not in st.session_state:
            st.session_state['df_combined'] = {}

        # Track if "Show PL Results" button is clicked
        if 'pl_show_results_clicked' not in st.session_state:
            st.session_state['pl_show_results_clicked'] = False

        # If the button is clicked, update the flag
        if st.button('Show PL Results', key='pl_show_results'):
            st.session_state['pl_show_results_clicked'] = True

        # Only process and display results if the button has been clicked
        if st.session_state['pl_show_results_clicked']:
            if peer_comparison_enabled:
                st.subheader("Peer Comparison")
                # Use the pre-fetched and processed NSE data from session state
                nse_data = st.session_state['nse_results_quarters']

                # Pivot the data to the desired format
                formatted_data = pivot_data_to_format(nse_data, selected_search_terms_quarters)

                # Retrieve sonata_data from session state (ensure it exists first)
                sonata_data = aggregate_data_single_row_merged(selected_search_terms_quarters, st.session_state['pl_data_frames'], selected_quarters)

                for term, df in formatted_data.items():
                    # If Basic (₹), skip convert_lakhs_to_crores
                    if term != 'Basic (₹)':
                        df = convert_lakhs_to_crores(df)  # Apply the lakhs to crores conversion for all terms except 'Basic (₹)'

                    df = df.reset_index()

                    # Filter peer data based on selected companies
                    df = df[df['Company'].isin(selected_companies)]

                    # Combine Sonata data and peers' data, removing duplicates
                    sonata_columns = [col for col in df.columns if col in sonata_data[term].columns and col != 'Company']
                    peer_columns = [col for col in df.columns if col not in sonata_data[term].columns and col != 'Company']

                    unique_columns = ['Company'] + sonata_columns + peer_columns

                    if term in sonata_data:
                        sonata_data[term] = sonata_data[term].reindex(columns=unique_columns, fill_value=0)
                    df = df.reindex(columns=unique_columns, fill_value=0)

                    # Concatenate the data
                    df_combined = pd.concat([sonata_data.get(term, pd.DataFrame()), df], axis=0, ignore_index=True)
                    st.session_state['df_combined'][term] = df_combined

                    # Display the combined DataFrame
                    st.dataframe(df_combined)

                    # Collect the DataFrame in the dictionary
                    term_data_dict[term] = df_combined

                    # Plot the data
                    cols = st.columns(1)
                    forecast_and_plot(df_combined, term, df_combined.columns[1:], cols[0], add_legend=False, chart_type='line')

            else:
                results_quarters = aggregate_data_single_row_merged(
                    selected_search_terms_quarters, 
                    st.session_state['pl_data_frames'], 
                    selected_quarters
                )

                st.session_state['pl_results_quarters'] = results_quarters
                st.session_state['sorted_pl_quarters'] = selected_quarters

                # Display the results without peer comparison
                display_consolidated_pl_results_quarters(
                    st.session_state['pl_results_quarters'], 
                    st.session_state['sorted_pl_quarters'], 
                    chart_type='area'
                )

                # Collect the DataFrame in the dictionary for non-peer comparison
                for term, df in st.session_state['pl_results_quarters'].items():
                    term_data_dict[term] = df

        # Show download button only if peer comparison is enabled and PowerPoint is generated
        if peer_comparison_enabled and term_data_dict:
            ppt_io = generate_ppt(term_data_dict, "P&L Presentation")
            st.download_button(
                label="Download P&L PowerPoint",
                data=ppt_io,
                file_name="profit_loss_visualization.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )

        

# Function to fetch and store NSE data upfront
def fetch_and_store_nse_data():
    symbols = ['TCS', 'INFY', 'HCLTECH','LTIM', 'WIPRO', 'COFORGE', 'PERSISTENT', 'MPHASIS', 'ZENSARTECH', 'BSOFT']
    all_data = {term: [] for term in ["Total Income", "Total expenses", "Net tax expense", "Profit for the year",
                                      "Revenue from operations", "Finance costs", "Depreciation and amortization expense",
                                      "Profit before exceptional item and tax", "Basic (₹)"]}

    for symbol in symbols:
        session = initialize_session(symbol)
        if session:
            financial_results = fetch_nse_financial_results(symbol, session)
            if financial_results:
                filtered_data = filter_data_by_year_and_consolidated(financial_results)
                links = create_financial_result_links(filtered_data)

                for idx, link in enumerate(links):
                    details = fetch_financial_result_details(link, session)
                    if details and "resultsData2" in details:
                        to_date = filtered_data[idx].get("toDate", "N/A")
                        quarter = get_quarter_from_date(to_date)
                        extracted_data = extract_required_fields(details["resultsData2"], quarter, symbol)
                        if extracted_data:
                            for key, value in extracted_data.items():
                                if key == "Quarter" or key == "Company":
                                    continue
                                all_data[key].append({"Quarter": extracted_data["Quarter"], "Company": symbol, "Value": value})

    # Convert to DataFrames and set index to 'Quarter' and 'Company'
    results_quarters = {}
    for key, data in all_data.items():
        if data:
            result_df = pd.DataFrame(data)
            if 'Quarter' in result_df.columns and 'Company' in result_df.columns:
                result_df.set_index(['Quarter', 'Company'], inplace=True)
            else:
                st.warning(f"No 'Quarter' or 'Company' column found in the data for {key}. Skipping this term.")
                continue

            results_quarters[key] = result_df

    st.session_state['nse_results_quarters'] = results_quarters
    st.session_state['nse_data_processed'] = True


def display_consolidated_pl_results_quarters(results, quarters, chart_type):
    valid_quarters = extract_and_sort_quarters(quarters)  # Ensure Quarters are sorted
    terms = list(results.keys())
    
    include_forecast = st.session_state.get('include_forecast_pl', False)

    first_graph = True
    for term in terms:
        #st.subheader(f"Results for {term}")
        df = results[term].reindex(columns=['Company'] + valid_quarters)
        st.dataframe(df)

        # Plotting
        if not df.empty:
            cols = st.columns(1)
            forecast_and_plot(df, term, valid_quarters, cols[0], add_legend=first_graph, chart_type=chart_type)
            first_graph = False

def kpi_page():
    kpi_data_dict = {}  # Dictionary to hold data for each KPI

    #st.subheader("Key Performance Indicators (KPIs)")

    # Track peer comparison state in session state
    if 'peer_comparison_enabled' not in st.session_state:
        st.session_state['peer_comparison_enabled'] = False

    peer_comparison_enabled = st.checkbox("Peer Comparison", value=st.session_state['peer_comparison_enabled'])
    st.session_state['peer_comparison_enabled'] = peer_comparison_enabled

    # Track if "Show KPI Results" button is clicked
    if 'kpi_show_results_clicked' not in st.session_state:
        st.session_state['kpi_show_results_clicked'] = False

    # If the button is clicked, update the flag
    if st.button('Show KPI Results', key='kpi_show_results'):
        st.session_state['kpi_show_results_clicked'] = True

    # Only process and display results if the button has been clicked or results already exist
    if st.session_state['kpi_show_results_clicked']:
        # Process Balance Sheet KPIs
        if 'merged_df1' in st.session_state:
            merged_df1 = st.session_state['merged_df1']

            # List of required terms for Balance Sheet
            terms_bs = [
                'Total assets', 
                'Total Equity', 
                'Total non-current liabilities', 
                'Total current liabilities'
            ]

            kpis_bs = {
                'Net Worth': pd.DataFrame(),
                'Net Assets': pd.DataFrame()
            }

            # Check if required data is available for Balance Sheet
            for term in terms_bs:
                if term not in merged_df1:
                    st.warning(f"Data for {term} not found in the combined data.")
                    return

            # Extract required data from Balance Sheet
            total_assets = merged_df1['Total assets']
            total_equity = merged_df1['Total Equity']
            total_non_current_liabilities = merged_df1['Total non-current liabilities']
            total_current_liabilities = merged_df1['Total current liabilities']

            companies = total_assets['Company'].unique() if peer_comparison_enabled else ['sonata-software']

            for company in companies:
                if not peer_comparison_enabled and company != 'sonata-software':
                    continue  # Skip other companies if peer comparison is disabled

                kpi_values_bs = {
                    'Net Worth': [],
                    'Net Assets': []
                }
                for col in total_assets.columns[1:]:
                    # Calculate KPIs for each financial year
                    total_assets_value = total_assets.loc[total_assets['Company'] == company, col].values[0] or 0
                    total_equity_value = total_equity.loc[total_equity['Company'] == company, col].values[0] or 0
                    total_non_current_liabilities_value = total_non_current_liabilities.loc[total_non_current_liabilities['Company'] == company, col].values[0] or 0
                    total_current_liabilities_value = total_current_liabilities.loc[total_current_liabilities['Company'] == company, col].values[0] or 0

                    net_worth_value = total_equity_value
                    net_assets_value = total_assets_value - total_non_current_liabilities_value - total_current_liabilities_value

                    # Append values to KPI dictionaries
                    kpi_values_bs['Net Worth'].append(net_worth_value)
                    kpi_values_bs['Net Assets'].append(net_assets_value)

                # Store KPI values in DataFrames
                for kpi in kpi_values_bs.keys():
                    kpis_bs[kpi][company] = kpi_values_bs[kpi]

            # Display KPIs for Balance Sheet
            for kpi_name, kpi_df in kpis_bs.items():
                st.subheader(f"KPI (Balance Sheet): {kpi_name}")
                kpi_df = pd.DataFrame(kpi_df).T  # Transpose to get financial years as columns
                kpi_df.columns = total_assets.columns[1:]  # Set financial year names as columns

                kpi_df.insert(0, 'Company', companies if peer_comparison_enabled else ['sonata-software'])

                kpi_df = kpi_df.reset_index(drop=True)
                st.dataframe(kpi_df)

                # Collect the DataFrame in the dictionary
                kpi_data_dict[kpi_name] = kpi_df

                # Plotting and Forecasting
                cols = st.columns(1)
                chart_type = 'line' if peer_comparison_enabled else 'area'
                forecast_and_plot(kpi_df, kpi_name, kpi_df.columns[1:], cols[0], chart_type=chart_type)
                st.markdown("**Note:** Data indicating zero may imply that the data is not present or not extracted due to improper headings.")
        else:
            st.warning("Balance Sheet data not found. Please process the files in the Balance Sheet section.")

        # Process Profit & Loss KPIs
        if 'df_combined' in st.session_state:
            df_combined = st.session_state['df_combined']

            # List of required terms for Profit & Loss
            terms_pl = [
                'Revenue from operations', 
                'Finance costs', 
                'Depreciation and amortization expense', 
                'Profit before exceptional item and tax', 
                'Basic (₹)', 
                'Profit for the year',
                'Net tax expense'
            ]

            kpis_pl = {
                'Net Sales': pd.DataFrame(),
                'EBITDA': pd.DataFrame(),
                'PAT': pd.DataFrame(),
                'EPS': pd.DataFrame(),
                'EBITDA Margin %': pd.DataFrame(),
                'Net Profit Margin %': pd.DataFrame(),
                'EBIT': pd.DataFrame()
            }

            # Check if required data is available for Profit & Loss
            for term in terms_pl:
                if term not in df_combined:
                    st.warning(f"Data for {term} not found in the combined data.")
                    return

            # Extract required data from Profit & Loss
            revenue = df_combined['Revenue from operations']
            finance_costs = df_combined['Finance costs']
            depreciation = df_combined['Depreciation and amortization expense']
            profit_before_tax = df_combined['Profit before exceptional item and tax']
            basic_eps = df_combined['Basic (₹)']
            profit_for_period = df_combined['Profit for the year']
            net_tax_expense = df_combined['Net tax expense']

            companies = revenue['Company'].unique() if peer_comparison_enabled else ['sonata-software']

            for company in companies:
                if not peer_comparison_enabled and company != 'sonata-software':
                    continue  # Skip other companies if peer comparison is disabled

                kpi_values_pl = {
                    'Net Sales': [],
                    'EBITDA': [],
                    'PAT': [],
                    'EPS': [],
                    'EBITDA Margin %': [],
                    'Net Profit Margin %': [],
                    'EBIT': []
                }
                for col in revenue.columns[1:]:
                    # Calculate KPIs for each quarter
                    net_sales_value = revenue.loc[revenue['Company'] == company, col].values[0]
                    profit_before_tax_value = profit_before_tax.loc[profit_before_tax['Company'] == company, col].values[0]
                    depreciation_value = depreciation.loc[depreciation['Company'] == company, col].values[0]
                    finance_costs_value = finance_costs.loc[finance_costs['Company'] == company, col].values[0]
                    pat_value = profit_for_period.loc[profit_for_period['Company'] == company, col].values[0]
                    eps_value = basic_eps.loc[basic_eps['Company'] == company, col].values[0]
                    net_tax_expense_value = net_tax_expense.loc[net_tax_expense['Company'] == company, col].values[0]

                    ebitda_value = profit_before_tax_value + depreciation_value + finance_costs_value
                    ebitda_margin_value = (ebitda_value / net_sales_value) * 100 if net_sales_value != 0 else 0
                    net_profit_margin_value = (pat_value / net_sales_value) * 100 if net_sales_value != 0 else 0
                    ebit_value = pat_value + net_tax_expense_value + finance_costs_value

                    # Append values to KPI dictionaries
                    kpi_values_pl['Net Sales'].append(net_sales_value)
                    kpi_values_pl['EBITDA'].append(ebitda_value)
                    kpi_values_pl['PAT'].append(pat_value)
                    kpi_values_pl['EPS'].append(eps_value)
                    kpi_values_pl['EBITDA Margin %'].append(ebitda_margin_value)
                    kpi_values_pl['Net Profit Margin %'].append(net_profit_margin_value)
                    kpi_values_pl['EBIT'].append(ebit_value)

                # Store KPI values in DataFrames
                for kpi in kpi_values_pl.keys():
                    kpis_pl[kpi][company] = kpi_values_pl[kpi]

            # Display KPIs for Profit & Loss
            for kpi_name, kpi_df in kpis_pl.items():
                st.subheader(f"KPI (Profit & Loss): {kpi_name}")
                kpi_df = pd.DataFrame(kpi_df).T  # Transpose to get quarters as columns
                kpi_df.columns = revenue.columns[1:]  # Set quarter names as columns

                # Insert the Company column correctly based on the number of rows in kpi_df
                kpi_df.insert(0, 'Company', companies if peer_comparison_enabled else ['sonata-software'])

                # Reset index to remove the index column
                kpi_df = kpi_df.reset_index(drop=True)

                # Display DataFrame
                st.dataframe(kpi_df)

                # Collect the DataFrame in the dictionary
                kpi_data_dict[kpi_name] = kpi_df

                # Plotting and Forecasting
                cols = st.columns(1)
                chart_type = 'line' if peer_comparison_enabled else 'area'
                forecast_and_plot(kpi_df, kpi_name, kpi_df.columns[1:], cols[0], chart_type=chart_type)
                st.markdown("**Note:** Data indicating zero may imply that the data is not present or not extracted due to improper headings.")
        else:
            st.warning("Profit & Loss data not found. Please process the files in the Profit and Loss section.")

        # PowerPoint Download for KPIs - ONLY if Peer Comparison is Enabled
        if peer_comparison_enabled and kpi_data_dict:
            ppt_io = generate_ppt(kpi_data_dict, "KPI Presentation")
            st.download_button(
                label="Download KPI PowerPoint",
                data=ppt_io,
                file_name="kpi_visualization.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )




# Load the selected page
if page == "Balance Sheet":
    balance_sheet_page()
elif page == "Profit & Loss":
    profit_loss_page()
elif page == "KPI":
    kpi_page()
